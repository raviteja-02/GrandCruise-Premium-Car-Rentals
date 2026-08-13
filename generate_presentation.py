from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# 16:9 aspect ratio in inches
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

slides_content = [
    {
        "title": "Project Title",
        "content": "GrandCruise: Modern Full-Stack Car Booking Platform"
    },
    {
        "title": "Abstract",
        "content": (
            "- GrandCruise is a modern full-stack car booking platform using Django and React.\n"
            "- Delivers a seamless and efficient experience for customers and administrators.\n"
            "- Backend APIs built with Django REST Framework.\n"
            "- Responsive frontend developed with React and TypeScript, styled using Tailwind CSS.\n"
            "- Supports JWT authentication for secure access.\n"
            "- Real-time vehicle availability and automated booking management.\n"
            "- Comprehensive admin tools for fleet management.\n"
            "- Modular architecture ensures scalability and maintainability.\n"
            "- Adheres to full-stack development best practices.\n"
            "- Sets a new benchmark in digital mobility solutions."
        )
    },
    {
        "title": "Introduction",
        "content": (
            "- The car rental industry is evolving with digital solutions.\n"
            "- GrandCruise addresses the need for a seamless, scalable, and user-friendly car booking platform.\n"
            "- Combines Django (backend) and React (frontend) for a robust, modern web application."
        )
    },
    {
        "title": "Literature Survey",
        "content": (
            "[1] Waspodo et al. (2011): Car rental management information systems.\n"
            "[2] Osman et al. (2017): Online car rental systems using web and SMS.\n"
            "[3] Fink & Reiners (2006): Short-term car rental logistics.\n"
            "[4] Khaled et al. (2015): Software requirements for online car rental.\n"
            "[5] Carroll & Grimes (1995): Product management in car rental industry.\n"
            "[6] Beck et al. (2001): Agile software development.\n"
            "[7] Abrahamsson et al. (2017): Agile methods review.\n"
            "[8] Soares & Moura (2015): Writing software requirements specifications."
        )
    },
    {
        "title": "Existing System pros and cons",
        "content": (
            "Pros:\n"
            "- Established platforms for car rental exist.\n"
            "- Some offer online booking and fleet management.\n\n"
            "Cons:\n"
            "- Many lack real-time availability updates.\n"
            "- Limited automation in booking and admin processes.\n"
            "- Often not mobile-friendly or scalable.\n"
            "- Poor integration with modern payment and notification systems."
        )
    },
    {
        "title": "6.1 Proposed Modules/Algorithms",
        "content": (
            "- User Management (registration, authentication, profiles)\n"
            "- Car Inventory Management (CRUD, categories, features)\n"
            "- Booking Management (create, update, cancel, confirm)\n"
            "- Admin Panel (fleet, bookings, users, settings)\n"
            "- Real-time availability and status updates\n"
            "- Filtering, search, and ordering for cars"
        )
    },
    {
        "title": "6.2 Dataset Description",
        "content": (
            "- User registrations and profiles\n"
            "- Vehicle inventory and attributes\n"
            "- Booking records and history\n"
            "- System-wide configuration values"
        )
    },
    {
        "title": "6.3 Proposed Work Architecture",
        "content": (
            "- Backend: Django, Django REST Framework, JWT Auth\n"
            "- Frontend: React, TypeScript, Tailwind CSS\n"
            "- Database: SQLite (default), Django ORM\n"
            "- APIs: RESTful, with filtering, search, and ordering\n"
            "- Deployment: Modular, scalable, maintainable"
        )
    },
    {
        "title": "Project Screenshots",
        "content": (
            "- Home Page and Car Details\n"
            "- User Registration and Login\n"
            "- Cars Listing and Car Details\n"
            "- Car Booking and User Bookings\n"
            "- Admin Bookings Overview and Car Management\n"
            "- System Settings Page\n"
            "- API Root Endpoints\n"
            "(Add screenshots here in your PPT manually)"
        )
    },
    {
        "title": "Conclusion",
        "content": (
            "GrandCruise demonstrates the implementation of a modern, full-stack car booking platform using Django and React. "
            "The platform provides a seamless user experience for both customers and administrators, with comprehensive backend APIs, responsive frontend interfaces, and advanced booking management capabilities. "
            "Its modular architecture ensures scalability and maintainability, setting new standards for user experience, system reliability, and business process automation in the premium mobility sector."
        )
    },
    {
        "title": "Future Scope",
        "content": (
            "- Mobile app development (React Native/Flutter)\n"
            "- In-app wallet or UPI-based payments\n"
            "- Real-time notifications (WebSockets, push)\n"
            "- Advanced analytics and AI-driven recommendations\n"
            "- Integration with GPS, insurance, and maintenance services\n"
            "- Blockchain contracts, IoT diagnostics, VR showrooms, ride-sharing integration"
        )
    },
    {
        "title": "References",
        "content": (
            "[1] Waspodo, Bayu, Qurrotul Aini, and Syamsuri Nur. \"Development of car rental management information system.\" In Proceeding International Conference on Information Systems For Business Competitiveness (ICISBC), pp. 101-105. 2011.\n"
            "[2] Osman, Mohd Nizam, Nurzaid Md Zain, Zulfikri Paidi, Khairul Anwar Sedek, Mohamad NajmuddinYusoff, and Mushahadah Maghribi. \"Online Car Rental System Using Web-Based and SMS Technology.\" Computing Research & Innovation (CRINN) 2 (2017): 277.\n"
            "[3] Fink, Andreas, and Torsten Reiners. \"Modeling and solving the short-term car rental logistics problem.\" Transportation Research Part E: Logistics and Transportation Review 42, no. 4 (2006): 272-292.\n"
            "[4] Khaled, Mr Shah Mostafa, Shamsil Arefin, Datta Sree Rajib Kumar, and Ariful Hossain Tuhin. \"Software Requirements Specification for Online Car Rental System.\" (2015).\n"
            "[5] Carroll, William J., and Richard C. Grimes. \"Evolutionary change in product management: Experiences in the car rental industry.\" Interfaces 25, no. 5 (1995): 84-104.\n"
            "[6] Beck, Kent, Mike Beedle, Arie Van Bennekum, Alistair Cockburn, Ward Cunningham, Martin Fowler, James Grenning et al. \"Manifesto for agile software development.\" (2001): 2006.\n"
            "[7] Abrahamsson, Pekka, Outi Salo, Jussi Ronkainen, and Juhani Warsta. \"Agile software development methods: Review and analysis.\" arXiv preprint arXiv:1709.08439 (2017).\n"
            "[8] Soares, Hécio A., and Raimundo S. Moura. \"A methodology to guide writing Software Requirements Specification document.\" In 2015 Latin American Computing Conference (CLEI), pp. 1-11. IEEE, 2015."
        )
    }
]

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
for slide in slides_content:
    slide_layout = prs.slide_layouts[1]  # Title and Content
    s = prs.slides.add_slide(slide_layout)
    # Set title font
    title_shape = s.shapes.title
    title_shape.text = slide["title"]
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Times New Roman'
            run.font.size = Pt(40)
    # Set content font
    content_shape = s.placeholders[1]
    content_shape.text = slide["content"]
    for paragraph in content_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Times New Roman'
            run.font.size = Pt(28)

prs.save("Project_Presentation.pptx")
print("Presentation created as Project_Presentation.pptx with Times New Roman font and 16:9 aspect ratio") 